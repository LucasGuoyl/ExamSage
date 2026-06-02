"""Ingest course materials (PDF / PPTX / Markdown / JSON) into normalized text.

Each ingest function returns a list of (text, metadata) tuples — one per
logical unit (PDF page, PPT slide, markdown section, JSON record).
Chunking happens in the next stage; this layer just extracts.
"""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Iterable

import fitz                              # PyMuPDF
from pptx import Presentation            # python-pptx


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def ingest_pdf(path: str | Path) -> list[tuple[str, dict]]:
    """Extract text from a PDF, one entry per page."""
    path = Path(path)
    doc = fitz.open(path)
    out: list[tuple[str, dict]] = []
    for page_idx, page in enumerate(doc):
        text = page.get_text("text").strip()
        if not text:
            continue
        out.append((text, {
            "source": path.name,
            "type": "pdf",
            "page": page_idx + 1,
            "total_pages": len(doc),
        }))
    doc.close()
    return out


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def ingest_pptx(path: str | Path) -> list[tuple[str, dict]]:
    """Extract text from a PPTX, one entry per slide.

    Slide text aggregates: title + all text frames + speaker notes.
    """
    path = Path(path)
    prs = Presentation(str(path))
    out: list[tuple[str, dict]] = []
    for slide_idx, slide in enumerate(prs.slides):
        parts: list[str] = []
        title = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None and shape == slide.shapes.title:
                title = text
            parts.append(text)
        # speaker notes (often hold "this will be on the exam" hints!)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[SPEAKER NOTES] {notes}")
        full = "\n\n".join(parts).strip()
        if not full:
            continue
        out.append((full, {
            "source": path.name,
            "type": "pptx",
            "slide": slide_idx + 1,
            "total_slides": len(prs.slides),
            "title": title,
            "has_notes": bool(notes),
        }))
    return out


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------
_HEADER_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$", re.MULTILINE)


def ingest_markdown(path: str | Path) -> list[tuple[str, dict]]:
    """Split a markdown file by headers; each section becomes one entry."""
    path = Path(path)
    raw = path.read_text(encoding="utf-8")
    sections: list[tuple[str, dict]] = []
    # split on H1/H2 boundaries
    parts = re.split(r"\n(?=#{1,3} )", raw)
    for i, part in enumerate(parts):
        part = part.strip()
        if not part:
            continue
        header_match = _HEADER_RE.match(part)
        header = header_match.group(2) if header_match else None
        sections.append((part, {
            "source": path.name,
            "type": "markdown",
            "section_idx": i,
            "header": header,
        }))
    return sections


# ---------------------------------------------------------------------------
# JSON past-paper format
# ---------------------------------------------------------------------------
def ingest_questions_json(path: str | Path) -> list[dict]:
    """Load past questions from a JSON file.

    Expected schema (flexible):
    [
      {"id": "2023_q1", "text": "...", "year": 2023, "answer": "...", "type": "MCQ"},
      ...
    ]
    """
    path = Path(path)
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, list):
        raise ValueError(f"{path}: expected a JSON list of question dicts")
    return data


# ---------------------------------------------------------------------------
# Directory walker
# ---------------------------------------------------------------------------
_SLIDE_EXTS = {".pdf", ".pptx", ".md"}
_QUESTION_EXTS = {".json", ".pdf", ".md"}


def ingest_directory(
    root: str | Path,
    subdir: str = "slides",
) -> list[tuple[str, dict]]:
    """Walk a directory and ingest all supported files in <root>/<subdir>."""
    root = Path(root) / subdir
    if not root.exists():
        return []
    all_units: list[tuple[str, dict]] = []
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        ext = path.suffix.lower()
        if ext == ".pdf":
            all_units.extend(ingest_pdf(path))
        elif ext == ".pptx":
            all_units.extend(ingest_pptx(path))
        elif ext == ".md":
            all_units.extend(ingest_markdown(path))
    return all_units


def discover_past_papers(root: str | Path) -> list[dict]:
    """Find all past-paper sources under <root>/past_papers/ and merge into one list.

    Supports .json (rich format) and .md / .pdf (will parse into rough questions).
    """
    root = Path(root) / "past_papers"
    if not root.exists():
        return []
    questions: list[dict] = []
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix.lower() == ".json":
            questions.extend(ingest_questions_json(path))
        elif path.suffix.lower() == ".md":
            # naive: each "##" header becomes a question
            md_units = ingest_markdown(path)
            for text, meta in md_units:
                if meta.get("header"):
                    questions.append({
                        "id": f"{path.stem}_{meta['section_idx']}",
                        "text": text,
                        "paper": path.stem,
                    })
        elif path.suffix.lower() == ".pdf":
            # naive: each page becomes one "question text"
            pdf_units = ingest_pdf(path)
            for text, meta in pdf_units:
                questions.append({
                    "id": f"{path.stem}_p{meta['page']}",
                    "text": text,
                    "paper": path.stem,
                })
    return questions
